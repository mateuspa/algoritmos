from pptx import Presentation
from pptx.util import Inches

# Create presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide 1: Introduction
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "About Me"
slide.placeholders[1].text = (
    "Name: [Your Name]\n"
    "Profession: Control and Automation Engineer\n"
    "Passionate about technology, innovation, and solving complex challenges."
)

# Slide 2: Academic Background
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Education & Continuous Learning"
slide.placeholders[1].text = (
    "- B.Sc. in Control and Automation Engineering – UNIFEI\n"
    "- Postgraduate in Marketing Management and Market Strategies – FDC\n"
    "- Currently pursuing a degree in Cross-Platform Software Development – FATEC"
)

# Slide 3: Technical and Strategic Blend
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Bridging Technology and Business"
slide.placeholders[1].text = (
    "- Strong foundation in automation, control systems, and process optimization\n"
    "- Proficient in software development for scalable solutions\n"
    "- Experience in aligning technical solutions with business strategy"
)

# Slide 4: Focus Areas
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Areas of Expertise"
slide.placeholders[1].text = (
    "- Industrial Automation\n"
    "- Software Development (Web, Mobile, Embedded)\n"
    "- Data Science & Analytics\n"
    "- Strategic Market Positioning"
)

# Slide 5: Why Me?
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "What I Bring to Your Consulting Firm"
slide.placeholders[1].text = (
    "- Multidisciplinary skill set\n"
    "- Problem-solving mindset\n"
    "- Eagerness to learn and adapt\n"
    "- Passion for delivering impactful solutions"
)

# Slide 6: Final Remarks
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Let’s Build the Future Together"
slide.placeholders[1].text = (
    "- Open to challenging projects\n"
    "- Committed to continuous growth\n"
    "- Ready to contribute from day one"
)

# Save presentation
prs.save("Consulting_Firm_Presentation.pptx")